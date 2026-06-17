from __future__ import annotations

import html
import os
import re
import shlex
import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any

from rag.bailian_ocr import BailianOcrClient, OcrResult
from rag.mineru_loader import MineruDocumentLoader
from rag.models import ParsedBlockDocument
from rag.parse_quality import QualitySignals, evaluate_parse_quality, merge_quality
from rag.summary_index import SummaryIndex
from schemas.rag import DocumentBlock, ParseQuality


TEXT_FILE_TYPES = {"txt", "text", "srt", "vtt"}
IMAGE_FILE_TYPES = {"png", "jpg", "jpeg", "webp"}
SUPPORTED_FILE_TYPES = {
    "pdf",
    "docx",
    "doc",
    "pptx",
    "ppt",
    "md",
    "markdown",
    "xlsx",
    "xls",
    *TEXT_FILE_TYPES,
    *IMAGE_FILE_TYPES,
}
ATOMIC_BLOCK_TYPES = {"table", "image", "chart", "formula", "code"}


class DocumentParserRouter:
    def __init__(
        self,
        mineru_loader: MineruDocumentLoader | None = None,
        ocr_client: BailianOcrClient | None = None,
    ) -> None:
        self.mineru_loader = mineru_loader or MineruDocumentLoader()
        self.ocr_client = ocr_client or BailianOcrClient.from_env()
        self.summary_index = SummaryIndex()

    def parse_bytes(
        self,
        *,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        document_type: str | None = None,
        content_type: str | None = None,
        source_path: str | None = None,
        high_precision: bool = False,
    ) -> ParsedBlockDocument:
        file_type = detect_file_type(filename, document_type, content_type)
        if file_type not in SUPPORTED_FILE_TYPES:
            file_type = "txt"

        try:
            if file_type == "pdf":
                blocks, quality, parser, warnings = self._parse_pdf(
                    content, filename, document_id, source_title, source_path, high_precision
                )
            elif file_type == "docx":
                blocks, quality, parser, warnings = self._parse_docx(
                    content, filename, document_id, source_title, source_path, high_precision
                )
            elif file_type == "doc":
                blocks, quality, parser, warnings = self._parse_legacy_office(
                    content, filename, document_id, source_title, source_path, "docx", high_precision
                )
            elif file_type == "pptx":
                blocks, quality, parser, warnings = self._parse_pptx(
                    content, filename, document_id, source_title, source_path, high_precision
                )
            elif file_type == "ppt":
                blocks, quality, parser, warnings = self._parse_legacy_office(
                    content, filename, document_id, source_title, source_path, "pptx", high_precision
                )
            elif file_type in {"srt", "vtt"}:
                text = decode_text(content)
                blocks = parse_transcript_blocks(
                    text=text,
                    document_id=document_id,
                    file_type=file_type,
                    source_title=source_title,
                    source_path=source_path,
                    parse_engine=f"{file_type}-subtitle-parser",
                )
                quality = evaluate_parse_quality(
                    QualitySignals(native_text_chars=len(text), paragraph_count=len(blocks)),
                    high_precision=False,
                )
                quality = mark_text_native_quality(quality)
                parser = f"{file_type}-subtitle-parser"
                warnings = []
            elif file_type in {"md", "markdown"}:
                text = decode_text(content)
                blocks = parse_markdown_blocks(
                    text=text,
                    document_id=document_id,
                    file_type="md",
                    source_title=source_title,
                    source_path=source_path,
                    parse_engine="markdown-it-py",
                )
                quality = evaluate_parse_quality(
                    QualitySignals(
                        native_text_chars=sum(len(block.contentText) for block in blocks),
                        paragraph_count=sum(1 for block in blocks if block.blockType in {"heading", "text", "list"}),
                        table_count=sum(1 for block in blocks if block.blockType == "table"),
                        image_count=sum(1 for block in blocks if block.blockType == "image"),
                    ),
                    high_precision=False,
                )
                quality = mark_text_native_quality(quality)
                parser = "markdown-it-py"
                warnings = []
            elif file_type in {"xlsx", "xls"}:
                blocks, quality, parser, warnings = self._parse_spreadsheet(
                    content, filename, document_id, source_title, source_path, file_type, high_precision
                )
            elif file_type in IMAGE_FILE_TYPES:
                blocks, quality, parser, warnings = self._parse_image(
                    content, filename, document_id, source_title, source_path, file_type
                )
            else:
                text = decode_text(content)
                if looks_like_transcript(text):
                    blocks = parse_transcript_blocks(
                        text=text,
                        document_id=document_id,
                        file_type=file_type,
                        source_title=source_title,
                        source_path=source_path,
                        parse_engine="timestamp-transcript-parser",
                    )
                    parser = "timestamp-transcript-parser"
                else:
                    blocks = parse_plain_text_blocks(
                        text=text,
                        document_id=document_id,
                        file_type=file_type,
                        source_title=source_title,
                        source_path=source_path,
                        parse_engine="text-encoding-detector",
                    )
                    parser = "text-encoding-detector"
                quality = evaluate_parse_quality(
                    QualitySignals(native_text_chars=len(text), paragraph_count=len(blocks)),
                    high_precision=False,
                )
                quality = mark_text_native_quality(quality)
                warnings = []
        except Exception as exc:
            fallback_text = decode_text(content)
            blocks = parse_plain_text_blocks(
                text=fallback_text,
                document_id=document_id,
                file_type=file_type,
                source_title=source_title,
                source_path=source_path,
                parse_engine="fallback-text",
            )
            quality = evaluate_parse_quality(QualitySignals(native_text_chars=len(fallback_text)), high_precision)
            parser = "fallback-text"
            warnings = [f"native parser failed: {exc}"]

        return self._finalize(document_id, blocks, parser, quality, warnings)

    def parse_text(
        self,
        *,
        document_id: str,
        title: str,
        document_type: str,
        source_path: str | None,
        content: str,
        parser: str,
    ) -> ParsedBlockDocument:
        file_type = normalize_file_type(document_type)
        if file_type in {"srt", "vtt"} or looks_like_transcript(content):
            blocks = parse_transcript_blocks(
                text=content,
                document_id=document_id,
                file_type=file_type,
                source_title=title,
                source_path=source_path,
                parse_engine=parser or "manual-transcript",
            )
        elif file_type in {"md", "markdown"}:
            blocks = parse_markdown_blocks(
                text=content,
                document_id=document_id,
                file_type="md",
                source_title=title,
                source_path=source_path,
                parse_engine=parser or "manual-markdown",
            )
        else:
            blocks = parse_plain_text_blocks(
                text=content,
                document_id=document_id,
                file_type=file_type,
                source_title=title,
                source_path=source_path,
                parse_engine=parser or "manual-text",
            )
        quality = evaluate_parse_quality(
            QualitySignals(native_text_chars=len(content), paragraph_count=len(blocks)),
            high_precision=False,
        )
        quality = mark_text_native_quality(quality)
        return self._finalize(document_id, blocks, parser or "manual-text", quality, [])

    def _parse_pdf(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        high_precision: bool,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        warnings: list[str] = []
        if self.mineru_loader.command:
            parsed = self.mineru_loader.load_bytes(content=content, filename=filename, content_type="application/pdf")
            if parsed.parser == "mineru" and parsed.text.strip():
                blocks = parse_markdown_blocks(
                    text=parsed.text,
                    document_id=document_id,
                    file_type="pdf",
                    source_title=source_title,
                    source_path=source_path,
                    parse_engine="mineru",
                )
                quality = evaluate_parse_quality(
                    QualitySignals(native_text_chars=len(parsed.text), paragraph_count=len(blocks)),
                    high_precision=high_precision,
                )
                return blocks, quality, "mineru", warnings
            warnings.append(f"MinerU did not return usable blocks: {parsed.parser}")

        blocks, native_warnings = parse_pdf_native_blocks(
            content=content,
            document_id=document_id,
            source_title=source_title,
            source_path=source_path,
            ocr_client=self.ocr_client,
        )
        warnings.extend(native_warnings)
        quality = evaluate_parse_quality(
            QualitySignals(
                native_text_chars=sum(len(block.contentText) for block in blocks),
                paragraph_count=len(blocks),
            ),
            high_precision=high_precision,
        )
        return blocks, quality, "pdf-native", warnings

    def _parse_docx(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        high_precision: bool,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        blocks, quality, warnings = parse_docx_blocks(
            content=content,
            document_id=document_id,
            source_title=source_title,
            source_path=source_path,
        )
        if high_precision:
            quality = quality.model_copy(
                update={
                    "highPrecision": True,
                    "needsSupplement": True,
                    "messages": [*quality.messages, "high precision parse requested"],
                }
            )
        parser = "python-docx"
        if quality.needsSupplement:
            supplement, supplement_quality, supplement_warnings = self._parse_office_pdf_supplement(
                content,
                filename,
                document_id,
                source_title,
                source_path,
                "docx",
            )
            if supplement:
                blocks.extend(mark_supplemental_blocks(supplement))
                quality = merge_quality(quality, supplement_quality)
                parser = "python-docx+libreoffice-pdf"
            else:
                warnings.extend(supplement_warnings)
        return blocks, quality, parser, warnings

    def _parse_pptx(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        high_precision: bool,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        blocks, quality, warnings = parse_pptx_blocks(
            content=content,
            document_id=document_id,
            source_title=source_title,
            source_path=source_path,
            high_precision=high_precision,
        )
        parser = "python-pptx"
        if quality.needsSupplement:
            supplement, supplement_quality, supplement_warnings = self._parse_office_pdf_supplement(
                content,
                filename,
                document_id,
                source_title,
                source_path,
                "pptx",
            )
            if supplement:
                blocks.extend(mark_supplemental_blocks(supplement))
                quality = merge_quality(quality, supplement_quality)
                parser = "python-pptx+libreoffice-pdf"
            else:
                warnings.extend(supplement_warnings)
        return blocks, quality, parser, warnings

    def _parse_legacy_office(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        native_target: str,
        high_precision: bool,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        warnings: list[str] = []
        with tempfile.TemporaryDirectory(prefix="rag-office-") as tmp:
            tmp_dir = Path(tmp)
            input_path = tmp_dir / filename
            input_path.write_bytes(content)
            native_path = convert_with_libreoffice(input_path, native_target, tmp_dir)
            pdf_path = convert_with_libreoffice(input_path, "pdf", tmp_dir)

            blocks: list[DocumentBlock] = []
            qualities: list[ParseQuality] = []
            parsers: list[str] = []
            if native_path and native_path.exists():
                native_content = native_path.read_bytes()
                if native_target == "docx":
                    native_blocks, native_quality, native_parser, native_warnings = self._parse_docx(
                        native_content,
                        native_path.name,
                        document_id,
                        source_title,
                        source_path,
                        high_precision,
                    )
                else:
                    native_blocks, native_quality, native_parser, native_warnings = self._parse_pptx(
                        native_content,
                        native_path.name,
                        document_id,
                        source_title,
                        source_path,
                        high_precision,
                    )
                blocks.extend(native_blocks)
                qualities.append(native_quality)
                parsers.append(f"libreoffice-{native_target}+{native_parser}")
                warnings.extend(native_warnings)
            else:
                warnings.append(f"LibreOffice could not convert {filename} to {native_target}")

            if pdf_path and pdf_path.exists():
                pdf_blocks, pdf_quality, pdf_parser, pdf_warnings = self._parse_pdf(
                    pdf_path.read_bytes(),
                    pdf_path.name,
                    document_id,
                    source_title,
                    source_path,
                    high_precision,
                )
                blocks.extend(mark_supplemental_blocks(pdf_blocks))
                qualities.append(pdf_quality)
                parsers.append(f"libreoffice-pdf+{pdf_parser}")
                warnings.extend(pdf_warnings)
            else:
                warnings.append(f"LibreOffice could not convert {filename} to pdf")

        quality = qualities[0] if qualities else evaluate_parse_quality(QualitySignals(), high_precision)
        for extra in qualities[1:]:
            quality = merge_quality(quality, extra)
        return blocks, quality, "+".join(parsers) if parsers else "libreoffice-unavailable", warnings

    def _parse_office_pdf_supplement(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        suffix: str,
    ) -> tuple[list[DocumentBlock], ParseQuality | None, list[str]]:
        warnings: list[str] = []
        with tempfile.TemporaryDirectory(prefix="rag-office-pdf-") as tmp:
            tmp_dir = Path(tmp)
            input_path = tmp_dir / filename
            if input_path.suffix.lower() != f".{suffix}":
                input_path = input_path.with_suffix(f".{suffix}")
            input_path.write_bytes(content)
            pdf_path = convert_with_libreoffice(input_path, "pdf", tmp_dir)
            if not pdf_path or not pdf_path.exists():
                warnings.append("LibreOffice PDF supplement is unavailable")
                return [], None, warnings
            blocks, quality, _parser, pdf_warnings = self._parse_pdf(
                pdf_path.read_bytes(),
                pdf_path.name,
                document_id,
                source_title,
                source_path,
                high_precision=False,
            )
            warnings.extend(pdf_warnings)
            return blocks, quality, warnings

    def _parse_spreadsheet(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        file_type: str,
        high_precision: bool,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        try:
            return parse_openpyxl_blocks(content, filename, document_id, source_title, source_path, high_precision)
        except Exception as exc:
            try:
                return parse_pandas_excel_blocks(content, document_id, source_title, source_path, file_type, exc)
            except Exception as fallback_exc:
                quality = evaluate_parse_quality(QualitySignals(), high_precision)
                return [], quality, "spreadsheet-unavailable", [
                    f"openpyxl parser failed: {exc}",
                    f"pandas parser failed: {fallback_exc}",
                ]

    def _parse_image(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        source_title: str,
        source_path: str | None,
        file_type: str,
    ) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
        warnings: list[str] = []
        ocr_result = self._ocr_image_bytes(content, filename=filename)
        warnings.extend(ocr_result.warnings)
        text = ocr_result.text
        parser = ocr_result.parser if text else "ocr-unavailable"
        confidence = ocr_result.confidence if text else 0.2
        if not text:
            try:
                from PIL import Image
                import pytesseract

                image = Image.open(BytesIO(content))
                text = pytesseract.image_to_string(image, lang=os.getenv("OCR_LANG", "chi_sim+eng")).strip()
                if text:
                    parser = "pytesseract"
                    confidence = 0.75
            except Exception as exc:
                parser = "ocr-unavailable"
                warnings.append(f"OCR unavailable: {exc}")

        block = DocumentBlock(
            documentId=document_id,
            blockId=f"{document_id}-image-1",
            fileType=file_type,
            blockType="image",
            contentText=text or "[图片] OCR 未获得可索引文字",
            assetPath=source_path,
            parseEngine=parser,
            confidence=confidence,
            sourceTitle=source_title,
            sourcePath=source_path,
            metadata={
                "filename": filename,
                "ocrTextChars": len(text),
                **ocr_result.metadata,
            },
        )
        quality = evaluate_parse_quality(
            QualitySignals(native_text_chars=len(text), image_count=1),
            high_precision=False,
        )
        if text:
            quality = quality.model_copy(update={"score": max(quality.score, confidence), "needsSupplement": False})
        return [block], quality, parser, warnings

    def _ocr_image_bytes(self, content: bytes, *, filename: str) -> OcrResult:
        if self.ocr_client.enabled:
            return self.ocr_client.recognize_image_bytes(image_bytes=content, filename=filename)
        return OcrResult(text="", parser="bailian-qwen-ocr-disabled")

    def _finalize(
        self,
        document_id: str,
        blocks: list[DocumentBlock],
        parser: str,
        quality: ParseQuality,
        warnings: list[str],
    ) -> ParsedBlockDocument:
        normalized_blocks = normalize_blocks(document_id, blocks)
        text_blocks = [block for block in normalized_blocks if block.contentText.strip()]
        if not text_blocks:
            status = "FAILED"
        elif warnings or quality.score < 0.68 or any(block.confidence < 0.4 for block in text_blocks):
            status = "PARTIAL"
        else:
            status = "READY"

        chunks_for_summary = [
            _summary_chunk(document_id, index, block)
            for index, block in enumerate(text_blocks)
        ]
        summaries = self.summary_index.build(chunks_for_summary) if chunks_for_summary else {
            "documentSummary": "",
            "sectionSummaries": {},
        }
        return ParsedBlockDocument(
            blocks=normalized_blocks,
            parser=parser,
            status=status,
            parse_quality=replace_parse_quality_messages(quality, warnings),
            document_summary=summaries["documentSummary"],
            section_summaries=summaries["sectionSummaries"],
            warnings=warnings,
        )


def detect_file_type(filename: str, document_type: str | None, content_type: str | None) -> str:
    normalized = normalize_file_type(document_type)
    if normalized in SUPPORTED_FILE_TYPES and normalized not in {"document", "upload"}:
        return "md" if normalized == "markdown" else normalized

    suffix = Path(filename or "").suffix.lower().lstrip(".")
    if suffix:
        return "md" if suffix == "markdown" else suffix
    if content_type and content_type.startswith("image/"):
        return content_type.split("/", 1)[1].lower()
    if content_type == "application/pdf":
        return "pdf"
    return "txt"


def normalize_file_type(document_type: str | None) -> str:
    if not document_type:
        return "txt"
    normalized = document_type.lower().strip().lstrip(".")
    if normalized == "markdown":
        return "md"
    if normalized == "text":
        return "txt"
    if normalized in {"subtitle", "subtitles", "caption", "captions"}:
        return "srt"
    if normalized in {"transcript", "asr", "video-transcript", "video_transcript"}:
        return "txt"
    return normalized


def decode_text(content: bytes) -> str:
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(content).best()
        if best is not None:
            return str(best)
    except Exception:
        pass
    try:
        import chardet

        detected = chardet.detect(content)
        encoding = detected.get("encoding")
        if encoding:
            return content.decode(encoding, errors="ignore")
    except Exception:
        pass
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


TIMESTAMP_RANGE_PATTERN = re.compile(
    r"(?P<start>(?:\d{1,2}:)?\d{1,2}:\d{2}(?:[,.]\d{1,3})?)\s*-->\s*"
    r"(?P<end>(?:\d{1,2}:)?\d{1,2}:\d{2}(?:[,.]\d{1,3})?)"
)
TIMESTAMP_PREFIX_PATTERN = re.compile(
    r"^\s*\[?(?P<start>(?:\d{1,2}:)?\d{1,2}:\d{2}(?:[,.]\d{1,3})?)\]?\s+(?P<text>.+)$"
)


def looks_like_transcript(text: str) -> bool:
    return bool(TIMESTAMP_RANGE_PATTERN.search(text) or any(TIMESTAMP_PREFIX_PATTERN.match(line) for line in text.splitlines()))


def parse_transcript_blocks(
    *,
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    normalized = text.replace("\ufeff", "").replace("\r\n", "\n").replace("\r", "\n")
    media_metadata = extract_transcript_media_metadata(normalized)
    cue_groups = [group.strip() for group in re.split(r"\n\s*\n", normalized) if group.strip()]
    for cue_index, group in enumerate(cue_groups, start=1):
        lines = [line.strip() for line in group.splitlines() if line.strip()]
        lines = [line for line in lines if line.upper() != "WEBVTT" and not line.isdigit()]
        if not lines:
            continue
        timestamp_index = next((index for index, line in enumerate(lines) if TIMESTAMP_RANGE_PATTERN.search(line)), -1)
        if timestamp_index >= 0:
            match = TIMESTAMP_RANGE_PATTERN.search(lines[timestamp_index])
            if not match:
                continue
            start_time = normalize_timestamp(match.group("start"))
            end_time = normalize_timestamp(match.group("end"))
            content = normalize_text("\n".join(lines[timestamp_index + 1 :]))
        else:
            prefix_match = TIMESTAMP_PREFIX_PATTERN.match(lines[0])
            if not prefix_match:
                continue
            start_time = normalize_timestamp(prefix_match.group("start"))
            end_time = None
            content = normalize_text(prefix_match.group("text") + "\n" + "\n".join(lines[1:]))
        if not content:
            continue
        section_title = f"{start_time} - {end_time}" if end_time else start_time
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-subtitle-{cue_index}",
                fileType=file_type,
                blockType="text",
                startTime=start_time,
                endTime=end_time,
                sectionTitle=section_title,
                contentText=content,
                parseEngine=parse_engine,
                confidence=0.95,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={
                    "cueIndex": cue_index,
                    "startTime": start_time,
                    "endTime": end_time,
                    "mediaType": "video",
                    "evidenceChannel": "subtitle",
                    **media_metadata,
                },
            )
        )
    if blocks:
        return blocks
    return parse_plain_text_blocks(
        text=text,
        document_id=document_id,
        file_type=file_type,
        source_title=source_title,
        source_path=source_path,
        parse_engine=parse_engine,
    )


def extract_transcript_media_metadata(text: str) -> dict[str, str]:
    """读取字幕或转写文本开头的可选视频播放地址。"""
    result: dict[str, str] = {}
    key_mapping = {
        "videourl": "videoUrl",
        "video_url": "videoUrl",
        "mediaurl": "mediaUrl",
        "media_url": "mediaUrl",
        "playbackurl": "playbackUrl",
        "playback_url": "playbackUrl",
    }
    for line in text.splitlines()[:12]:
        if ":" not in line:
            continue
        raw_key, raw_value = line.split(":", 1)
        key = raw_key.strip().lower()
        value = raw_value.strip()
        target_key = key_mapping.get(key)
        if target_key and value:
            result[target_key] = value
    return result


def normalize_timestamp(value: str) -> str:
    cleaned = value.replace(",", ".").split(".", 1)[0]
    parts = [int(part) for part in cleaned.split(":")]
    if len(parts) == 2:
        hours = 0
        minutes, seconds = parts
    else:
        hours, minutes, seconds = parts[-3:]
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}"


def parse_plain_text_blocks(
    *,
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    section_title: str | None = None
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", normalize_text(text)) if part.strip()]
    if not paragraphs and text.strip():
        paragraphs = [normalize_text(text)]
    for index, paragraph in enumerate(paragraphs, start=1):
        block_type = "heading" if looks_like_heading(paragraph) else "text"
        if block_type == "heading":
            section_title = paragraph.strip("# :：")
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-txt-{index}",
                fileType=file_type,
                blockType=block_type,
                sectionTitle=section_title,
                contentText=paragraph,
                parseEngine=parse_engine,
                confidence=0.92,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={"paragraphIndex": index},
            )
        )
    return blocks


def parse_markdown_blocks(
    *,
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
) -> list[DocumentBlock]:
    try:
        return parse_markdown_it_blocks(text, document_id, file_type, source_title, source_path, parse_engine)
    except Exception:
        fallback = parse_markdown_line_blocks(text, document_id, file_type, source_title, source_path, parse_engine)
        table_blocks = parse_markdown_tables(
            text,
            document_id,
            file_type,
            source_title,
            source_path,
            parse_engine,
            len(fallback),
            next((block.sectionTitle for block in reversed(fallback) if block.sectionTitle), None),
        )
        return fallback + table_blocks


def parse_markdown_it_blocks(
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
) -> list[DocumentBlock]:
    from markdown_it import MarkdownIt

    tokens = MarkdownIt("commonmark").parse(text)
    blocks: list[DocumentBlock] = []
    section_title: str | None = None
    index = 0
    cursor = 0
    while cursor < len(tokens):
        token = tokens[cursor]
        if token.type == "heading_open" and cursor + 1 < len(tokens):
            inline = tokens[cursor + 1]
            content = normalize_text(inline.content)
            if content:
                section_title = content
                index += 1
                blocks.append(
                    make_block(document_id, index, file_type, "heading", content, source_title, source_path, parse_engine, section_title)
                )
            cursor += 3
            continue
        if token.type == "paragraph_open" and cursor + 1 < len(tokens):
            inline = tokens[cursor + 1]
            image_children = [child for child in inline.children or [] if child.type == "image"]
            content = normalize_text(inline.content)
            for child in image_children:
                index += 1
                blocks.append(
                    make_block(
                        document_id,
                        index,
                        file_type,
                        "image",
                        child.content or child.attrs.get("alt", "") if child.attrs else "Markdown 图片",
                        source_title,
                        source_path,
                        parse_engine,
                        section_title,
                        asset_path=child.attrs.get("src") if child.attrs else None,
                    )
                )
            if content:
                index += 1
                blocks.append(
                    make_block(document_id, index, file_type, "text", content, source_title, source_path, parse_engine, section_title)
                )
            cursor += 3
            continue
        if token.type == "fence":
            index += 1
            blocks.append(
                make_block(
                    document_id,
                    index,
                    file_type,
                    "code",
                    token.content.strip(),
                    source_title,
                    source_path,
                    parse_engine,
                    section_title,
                    metadata={"language": token.info.strip()},
                )
            )
            cursor += 1
            continue
        if token.type in {"bullet_list_open", "ordered_list_open"}:
            list_lines: list[str] = []
            ordered = token.type == "ordered_list_open"
            cursor += 1
            while cursor < len(tokens) and tokens[cursor].type not in {"bullet_list_close", "ordered_list_close"}:
                if tokens[cursor].type == "inline" and tokens[cursor].content.strip():
                    prefix = f"{len(list_lines) + 1}." if ordered else "-"
                    list_lines.append(f"{prefix} {tokens[cursor].content.strip()}")
                cursor += 1
            if list_lines:
                index += 1
                blocks.append(
                    make_block(
                        document_id,
                        index,
                        file_type,
                        "list",
                        "\n".join(list_lines),
                        source_title,
                        source_path,
                        parse_engine,
                        section_title,
                    )
                )
            cursor += 1
            continue
        cursor += 1

    table_blocks = parse_markdown_tables(text, document_id, file_type, source_title, source_path, parse_engine, len(blocks), section_title)
    return blocks + table_blocks if blocks or table_blocks else parse_markdown_line_blocks(
        text, document_id, file_type, source_title, source_path, parse_engine
    )


def parse_markdown_line_blocks(
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    section_title: str | None = None
    buffer: list[str] = []
    block_index = 0
    in_code = False
    code_buffer: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if line.strip().startswith("```"):
            if in_code:
                block_index += 1
                blocks.append(
                    make_block(
                        document_id,
                        block_index,
                        file_type,
                        "code",
                        "\n".join(code_buffer),
                        source_title,
                        source_path,
                        parse_engine,
                        section_title,
                    )
                )
                code_buffer = []
                in_code = False
            else:
                if buffer:
                    block_index += 1
                    blocks.append(
                        make_block(
                            document_id,
                            block_index,
                            file_type,
                            "text",
                            "\n".join(buffer),
                            source_title,
                            source_path,
                            parse_engine,
                            section_title,
                        )
                    )
                    buffer = []
                in_code = True
            continue
        if in_code:
            code_buffer.append(line)
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        if heading:
            if buffer:
                block_index += 1
                blocks.append(
                    make_block(
                        document_id,
                        block_index,
                        file_type,
                        "text",
                        "\n".join(buffer),
                        source_title,
                        source_path,
                        parse_engine,
                        section_title,
                    )
                )
                buffer = []
            section_title = heading.group(2).strip()
            block_index += 1
            blocks.append(
                make_block(
                    document_id,
                    block_index,
                    file_type,
                    "heading",
                    section_title,
                    source_title,
                    source_path,
                    parse_engine,
                    section_title,
                )
            )
            continue
        if line.strip():
            buffer.append(line)
        elif buffer:
            block_index += 1
            text_block = "\n".join(buffer)
            block_type = "list" if re.match(r"^\s*(-|\d+\.)\s+", text_block) else "text"
            blocks.append(
                make_block(
                    document_id,
                    block_index,
                    file_type,
                    block_type,
                    text_block,
                    source_title,
                    source_path,
                    parse_engine,
                    section_title,
                )
            )
            buffer = []
    if buffer:
        block_index += 1
        text_block = "\n".join(buffer)
        blocks.append(
            make_block(
                document_id,
                block_index,
                file_type,
                "list" if re.match(r"^\s*(-|\d+\.)\s+", text_block) else "text",
                text_block,
                source_title,
                source_path,
                parse_engine,
                section_title,
            )
        )
    return blocks


def parse_markdown_tables(
    text: str,
    document_id: str,
    file_type: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
    start_index: int,
    section_title: str | None,
) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    lines = text.splitlines()
    cursor = 0
    index = start_index
    while cursor < len(lines) - 1:
        if "|" in lines[cursor] and re.match(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$", lines[cursor + 1]):
            table_lines = [lines[cursor], lines[cursor + 1]]
            cursor += 2
            while cursor < len(lines) and "|" in lines[cursor] and lines[cursor].strip():
                table_lines.append(lines[cursor])
                cursor += 1
            rows = [[cell.strip() for cell in row.strip().strip("|").split("|")] for row in table_lines if "---" not in row]
            index += 1
            blocks.append(
                DocumentBlock(
                    documentId=document_id,
                    blockId=f"{document_id}-md-table-{index}",
                    fileType=file_type,
                    blockType="table",
                    sectionTitle=section_title,
                    contentText=table_to_text(rows),
                    contentHtml=table_to_html(rows),
                    parseEngine=parse_engine,
                    confidence=0.88,
                    sourceTitle=source_title,
                    sourcePath=source_path,
                    metadata={"markdownTable": True},
                )
            )
            continue
        cursor += 1
    return blocks


def parse_pdf_native_blocks(
    *,
    content: bytes,
    document_id: str,
    source_title: str,
    source_path: str | None,
    ocr_client: BailianOcrClient | None = None,
) -> tuple[list[DocumentBlock], list[str]]:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)
    try:
        blocks = parse_pdf_with_pymupdf(tmp_path, document_id, source_title, source_path)
        if blocks:
            return blocks, []
        blocks = parse_pdf_with_pdfplumber(tmp_path, document_id, source_title, source_path)
        if blocks:
            return blocks, []
        blocks = parse_pdf_with_pypdf(tmp_path, document_id, source_title, source_path)
        if blocks:
            return blocks, []
        return parse_pdf_with_ocr(tmp_path, document_id, source_title, source_path, ocr_client=ocr_client)
    finally:
        tmp_path.unlink(missing_ok=True)


def parse_pdf_with_pymupdf(path: Path, document_id: str, source_title: str, source_path: str | None) -> list[DocumentBlock]:
    try:
        import fitz
    except Exception:
        return []
    blocks: list[DocumentBlock] = []
    with fitz.open(str(path)) as doc:
        for page_number, page in enumerate(doc, start=1):
            text = normalize_text(page.get_text("text") or "")
            if text:
                blocks.append(
                    make_block(
                        document_id,
                        page_number,
                        "pdf",
                        "text",
                        text,
                        source_title,
                        source_path,
                        "pymupdf",
                        None,
                        page_index=page_number,
                    )
                )
    return blocks


def parse_pdf_with_pdfplumber(path: Path, document_id: str, source_title: str, source_path: str | None) -> list[DocumentBlock]:
    try:
        import pdfplumber
    except Exception:
        return []
    blocks: list[DocumentBlock] = []
    with pdfplumber.open(str(path)) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            text = normalize_text(page.extract_text() or "")
            if text:
                blocks.append(
                    make_block(
                        document_id,
                        page_number,
                        "pdf",
                        "text",
                        text,
                        source_title,
                        source_path,
                        "pdfplumber",
                        None,
                        page_index=page_number,
                    )
                )
    return blocks


def parse_pdf_with_pypdf(path: Path, document_id: str, source_title: str, source_path: str | None) -> list[DocumentBlock]:
    try:
        from pypdf import PdfReader
    except Exception:
        return []
    reader = PdfReader(str(path))
    blocks: list[DocumentBlock] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text:
            blocks.append(
                make_block(
                    document_id,
                    page_number,
                    "pdf",
                    "text",
                    text,
                    source_title,
                    source_path,
                    "pypdf",
                    None,
                    page_index=page_number,
                )
            )
    return blocks


def parse_pdf_with_ocr(
    path: Path,
    document_id: str,
    source_title: str,
    source_path: str | None,
    ocr_client: BailianOcrClient | None = None,
) -> tuple[list[DocumentBlock], list[str]]:
    try:
        import fitz
    except Exception as exc:
        return [], [f"PDF OCR unavailable: {exc}"]

    warnings: list[str] = []
    blocks: list[DocumentBlock] = []
    with fitz.open(str(path)) as doc:
        for page_number, page in enumerate(doc, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            png_bytes = pixmap.tobytes("png")
            text = ""
            parser = "pymupdf-ocr"
            confidence = 0.2
            metadata: dict[str, Any] = {"pageIndex": page_number}

            if ocr_client is not None and ocr_client.enabled:
                result = ocr_client.recognize_image_bytes(
                    image_bytes=png_bytes,
                    filename=f"{path.stem}-page-{page_number}.png",
                    mime_type="image/png",
                )
                warnings.extend(result.warnings)
                text = result.text
                parser = result.parser if text else parser
                confidence = result.confidence if text else confidence
                metadata.update(result.metadata)

            if not text:
                local_text, local_warnings = _ocr_image_bytes_with_pytesseract(png_bytes, page_number=page_number)
                warnings.extend(local_warnings)
                if local_text.strip():
                    text = local_text
                    parser = "pymupdf-ocr"
                    confidence = 0.78

            text = normalize_text(text or "")
            if text:
                blocks.append(
                    make_block(
                        document_id,
                        page_number,
                        "pdf",
                        "text",
                        text,
                        source_title,
                        source_path,
                        parser,
                        None,
                        page_index=page_number,
                        metadata=metadata,
                    ).model_copy(update={"confidence": confidence})
                )

    return blocks, warnings


def _ocr_image_bytes_with_pytesseract(image_bytes: bytes, *, page_number: int | None = None) -> tuple[str, list[str]]:
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(BytesIO(image_bytes))
        return normalize_text(pytesseract.image_to_string(image, lang=os.getenv("OCR_LANG", "chi_sim+eng")) or ""), []
    except Exception as exc:
        location = f" page {page_number}" if page_number is not None else ""
        return "", [f"Local OCR unavailable{location}: {exc}"]


def parse_docx_blocks(
    *,
    content: bytes,
    document_id: str,
    source_title: str,
    source_path: str | None,
) -> tuple[list[DocumentBlock], ParseQuality, list[str]]:
    from docx import Document

    warnings: list[str] = []
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)
    try:
        document = Document(str(tmp_path))
    finally:
        tmp_path.unlink(missing_ok=True)

    blocks: list[DocumentBlock] = []
    section_title: str | None = None
    text_chars = 0
    paragraph_count = 0
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = normalize_text(paragraph.text)
        if not text:
            continue
        paragraph_count += 1
        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        block_type = "heading" if style_name.lower().startswith("heading") or looks_like_heading(text) else "text"
        if block_type == "heading":
            section_title = text.strip("# :：")
        text_chars += len(text)
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-docx-p{index}",
                fileType="docx",
                blockType=block_type,
                sectionTitle=section_title,
                contentText=text,
                parseEngine="python-docx",
                confidence=0.9,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={"paragraphIndex": index, "style": style_name},
            )
        )

    for table_index, table in enumerate(document.tables, start=1):
        rows = [[normalize_text(cell.text) for cell in row.cells] for row in table.rows]
        text = table_to_text(rows)
        if not text.strip():
            continue
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-docx-t{table_index}",
                fileType="docx",
                blockType="table",
                sectionTitle=section_title,
                contentText=text,
                contentHtml=table_to_html(rows),
                parseEngine="python-docx",
                confidence=0.86,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={"tableIndex": table_index},
            )
        )
        text_chars += len(text)

    image_count = len(document.inline_shapes)
    for image_index, shape in enumerate(document.inline_shapes, start=1):
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-docx-img{image_index}",
                fileType="docx",
                blockType="image",
                sectionTitle=section_title,
                contentText=f"图片 {image_index}",
                parseEngine="python-docx",
                confidence=0.35,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={
                    "imageIndex": image_index,
                    "width": safe_int(getattr(shape, "width", None)),
                    "height": safe_int(getattr(shape, "height", None)),
                },
            )
        )

    embedded_object_count = count_docx_embedded_objects(document)
    quality = evaluate_parse_quality(
        QualitySignals(
            native_text_chars=text_chars,
            paragraph_count=paragraph_count,
            table_count=len(document.tables),
            image_count=image_count,
            drawing_count=image_count,
            embedded_object_count=embedded_object_count,
        )
    )
    return blocks, quality, warnings


def parse_pptx_blocks(
    *,
    content: bytes,
    document_id: str,
    source_title: str,
    source_path: str | None,
    high_precision: bool,
) -> tuple[list[DocumentBlock], ParseQuality, list[str]]:
    from pptx import Presentation

    warnings: list[str] = []
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)
    try:
        presentation = Presentation(str(tmp_path))
    finally:
        tmp_path.unlink(missing_ok=True)

    blocks: list[DocumentBlock] = []
    text_chars = 0
    shape_count = 0
    table_count = 0
    image_count = 0
    text_box_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        section_title = extract_slide_title(slide)
        if section_title:
            blocks.append(
                DocumentBlock(
                    documentId=document_id,
                    blockId=f"{document_id}-pptx-s{slide_index}-title",
                    fileType="pptx",
                    blockType="heading",
                    slideIndex=slide_index,
                    sectionTitle=section_title,
                    contentText=section_title,
                    parseEngine="python-pptx",
                    confidence=0.9,
                    sourceTitle=source_title,
                    sourcePath=source_path,
                    metadata={"slideIndex": slide_index, "role": "title"},
                )
            )
            text_chars += len(section_title)
        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
                rows = [
                    [normalize_text(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                text = table_to_text(rows)
                if text:
                    blocks.append(
                        DocumentBlock(
                            documentId=document_id,
                            blockId=f"{document_id}-pptx-s{slide_index}-t{shape_index}",
                            fileType="pptx",
                            blockType="table",
                            slideIndex=slide_index,
                            sectionTitle=section_title,
                            contentText=text,
                            contentHtml=table_to_html(rows),
                            parseEngine="python-pptx",
                            confidence=0.84,
                            sourceTitle=source_title,
                            sourcePath=source_path,
                            metadata={"slideIndex": slide_index, "shapeIndex": shape_index},
                        )
                    )
                    text_chars += len(text)
                continue
            if getattr(shape, "has_text_frame", False):
                text = normalize_text("\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs))
                if text and text != section_title:
                    text_box_count += 1
                    blocks.append(
                        DocumentBlock(
                            documentId=document_id,
                            blockId=f"{document_id}-pptx-s{slide_index}-p{shape_index}",
                            fileType="pptx",
                            blockType="text",
                            slideIndex=slide_index,
                            sectionTitle=section_title,
                            contentText=text,
                            parseEngine="python-pptx",
                            confidence=0.86,
                            sourceTitle=source_title,
                            sourcePath=source_path,
                            metadata={"slideIndex": slide_index, "shapeIndex": shape_index},
                        )
                    )
                    text_chars += len(text)
            if hasattr(shape, "image"):
                image_count += 1
                blocks.append(
                    DocumentBlock(
                        documentId=document_id,
                        blockId=f"{document_id}-pptx-s{slide_index}-img{shape_index}",
                        fileType="pptx",
                        blockType="image",
                        slideIndex=slide_index,
                        sectionTitle=section_title,
                        contentText=f"幻灯片 {slide_index} 图片 {image_count}",
                        parseEngine="python-pptx",
                        confidence=0.32,
                        sourceTitle=source_title,
                        sourcePath=source_path,
                        metadata={"slideIndex": slide_index, "shapeIndex": shape_index},
                    )
                )
        notes_text = extract_slide_notes(slide)
        if notes_text:
            blocks.append(
                DocumentBlock(
                    documentId=document_id,
                    blockId=f"{document_id}-pptx-s{slide_index}-notes",
                    fileType="pptx",
                    blockType="text",
                    slideIndex=slide_index,
                    sectionTitle=section_title,
                    contentText=notes_text,
                    parseEngine="python-pptx",
                    confidence=0.82,
                    sourceTitle=source_title,
                    sourcePath=source_path,
                    metadata={"slideIndex": slide_index, "role": "notes"},
                )
            )
            text_chars += len(notes_text)

    quality = evaluate_parse_quality(
        QualitySignals(
            native_text_chars=text_chars,
            paragraph_count=text_box_count,
            table_count=table_count,
            image_count=image_count,
            shape_count=shape_count,
            text_box_count=text_box_count,
            drawing_count=image_count,
        ),
        high_precision=high_precision,
    )
    return blocks, quality, warnings


def parse_openpyxl_blocks(
    content: bytes,
    filename: str,
    document_id: str,
    source_title: str,
    source_path: str | None,
    high_precision: bool,
) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
    from openpyxl import load_workbook

    with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix or ".xlsx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)
    try:
        workbook = load_workbook(str(tmp_path), data_only=False)
    finally:
        tmp_path.unlink(missing_ok=True)

    blocks: list[DocumentBlock] = []
    total_cells = 0
    empty_cells = 0
    merged_count = 0
    text_chars = 0
    formula_count = 0
    for sheet in workbook.worksheets:
        cell_range = sheet.calculate_dimension()
        rows: list[list[str]] = []
        for row in sheet.iter_rows():
            rendered_row: list[str] = []
            for cell in row:
                total_cells += 1
                value = cell.value
                if value is None:
                    empty_cells += 1
                    rendered_row.append("")
                    continue
                rendered = str(value)
                if rendered.startswith("="):
                    formula_count += 1
                rendered_row.append(rendered)
            rows.append(rendered_row)
        merged_ranges = [str(item) for item in sheet.merged_cells.ranges]
        merged_count += len(merged_ranges)
        text = table_to_text(rows)
        if not text.strip():
            continue
        text_chars += len(text)
        blocks.append(
            DocumentBlock(
                documentId=document_id,
                blockId=f"{document_id}-xlsx-{sheet.title}",
                fileType="xlsx",
                blockType="table",
                sheetName=sheet.title,
                cellRange=cell_range,
                contentText=text,
                contentHtml=table_to_html(rows),
                parseEngine="openpyxl",
                confidence=0.84,
                sourceTitle=source_title,
                sourcePath=source_path,
                metadata={
                    "sheetName": sheet.title,
                    "cellRange": cell_range,
                    "mergedCells": merged_ranges,
                    "formulaCount": formula_count,
                },
            )
        )
    quality = evaluate_parse_quality(
        QualitySignals(
            native_text_chars=text_chars,
            table_count=len(blocks),
            merged_cell_count=merged_count,
            empty_cell_ratio=empty_cells / max(total_cells, 1),
        ),
        high_precision=high_precision,
    )
    return blocks, quality, "openpyxl", []


def parse_pandas_excel_blocks(
    content: bytes,
    document_id: str,
    source_title: str,
    source_path: str | None,
    file_type: str,
    original_error: Exception,
) -> tuple[list[DocumentBlock], ParseQuality, str, list[str]]:
    import pandas as pd

    sheets = pd.read_excel(BytesIO(content), sheet_name=None, header=None)
    blocks: list[DocumentBlock] = []
    text_chars = 0
    total_cells = 0
    empty_cells = 0
    for sheet_name, frame in sheets.items():
        rows = frame.fillna("").astype(str).values.tolist()
        total_cells += int(frame.size)
        empty_cells += int(frame.isna().sum().sum())
        text = table_to_text(rows)
        if text.strip():
            text_chars += len(text)
            blocks.append(
                DocumentBlock(
                    documentId=document_id,
                    blockId=f"{document_id}-{file_type}-{sheet_name}",
                    fileType=file_type,
                    blockType="table",
                    sheetName=str(sheet_name),
                    contentText=text,
                    contentHtml=table_to_html(rows),
                    parseEngine="pandas",
                    confidence=0.72,
                    sourceTitle=source_title,
                    sourcePath=source_path,
                    metadata={"sheetName": str(sheet_name), "fallbackReason": str(original_error)},
                )
            )
    quality = evaluate_parse_quality(
        QualitySignals(
            native_text_chars=text_chars,
            table_count=len(blocks),
            empty_cell_ratio=empty_cells / max(total_cells, 1),
        )
    )
    return blocks, quality, "pandas", [f"openpyxl parser failed: {original_error}"]


def convert_with_libreoffice(input_path: Path, target_ext: str, output_dir: Path) -> Path | None:
    command = os.getenv("LIBREOFFICE_COMMAND") or os.getenv("SOFFICE_COMMAND")
    executable = command or shutil.which("soffice") or shutil.which("libreoffice")
    if not executable:
        return None

    if command and ("{input}" in command or "{output}" in command or "{format}" in command):
        args = shlex.split(command.format(input=str(input_path), output=str(output_dir), format=target_ext))
    else:
        args = shlex.split(executable) + [
            "--headless",
            "--convert-to",
            target_ext,
            "--outdir",
            str(output_dir),
            str(input_path),
        ]
    try:
        subprocess.run(args, check=True, timeout=120, capture_output=True, text=True)
    except Exception:
        return None

    expected = output_dir / f"{input_path.stem}.{target_ext}"
    if expected.exists():
        return expected
    matches = list(output_dir.glob(f"*.{target_ext}"))
    return matches[0] if matches else None


def mark_supplemental_blocks(blocks: list[DocumentBlock]) -> list[DocumentBlock]:
    result = []
    for block in blocks:
        metadata = {**block.metadata, "supplemental": True}
        result.append(block.model_copy(update={"metadata": metadata, "confidence": min(block.confidence, 0.78)}))
    return result


def normalize_blocks(document_id: str, blocks: list[DocumentBlock]) -> list[DocumentBlock]:
    unique: list[DocumentBlock] = []
    seen: set[tuple[str, str, str, str | None, int | None, int | None, str | None]] = set()
    for block in blocks:
        text = normalize_text(block.contentText)
        if not text:
            continue
        signature = (
            block.blockType,
            text[:180],
            block.parseEngine,
            block.sectionTitle,
            block.pageIndex,
            block.slideIndex,
            block.sheetName,
        )
        if signature in seen:
            continue
        seen.add(signature)
        unique.append(block.model_copy(update={"contentText": text}))
    normalized: list[DocumentBlock] = []
    for index, block in enumerate(unique, start=1):
        normalized.append(block.model_copy(update={"blockId": f"{document_id}-b{index:04d}"}))
    return normalized


def replace_parse_quality_messages(quality: ParseQuality, warnings: list[str]) -> ParseQuality:
    messages = list(dict.fromkeys([*quality.messages, *warnings]))
    return quality.model_copy(update={"messages": messages})


def mark_text_native_quality(quality: ParseQuality) -> ParseQuality:
    return quality.model_copy(update={"score": max(quality.score, 0.9), "needsSupplement": False})


def _summary_chunk(document_id: str, index: int, block: DocumentBlock):
    from rag.models import Chunk

    return Chunk(
        chunk_id=f"{document_id}-summary-{index}",
        document_id=document_id,
        text=block.contentText,
        metadata={"sectionName": block.sectionTitle or "全文"},
    )


def make_block(
    document_id: str,
    index: int,
    file_type: str,
    block_type: str,
    content_text: str,
    source_title: str,
    source_path: str | None,
    parse_engine: str,
    section_title: str | None,
    *,
    page_index: int | None = None,
    asset_path: str | None = None,
    metadata: dict[str, Any] | None = None,
) -> DocumentBlock:
    return DocumentBlock(
        documentId=document_id,
        blockId=f"{document_id}-{file_type}-{index}",
        fileType=file_type,
        blockType=block_type,  # type: ignore[arg-type]
        pageIndex=page_index,
        sectionTitle=section_title,
        contentText=normalize_text(content_text),
        assetPath=asset_path,
        parseEngine=parse_engine,
        confidence=0.9,
        sourceTitle=source_title,
        sourcePath=source_path,
        metadata=metadata or {},
    )


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def looks_like_heading(text: str) -> bool:
    stripped = text.strip()
    if stripped.startswith("#"):
        return True
    return len(stripped) <= 36 and stripped.endswith((":", "："))


def table_to_text(rows: list[list[Any]]) -> str:
    lines = []
    for row in rows:
        rendered = [normalize_text(str(cell)) for cell in row]
        if any(rendered):
            lines.append("\t".join(rendered))
    return "\n".join(lines)


def table_to_html(rows: list[list[Any]]) -> str:
    html_rows = []
    for row in rows:
        cells = "".join(f"<td>{html.escape(normalize_text(str(cell)))}</td>" for cell in row)
        html_rows.append(f"<tr>{cells}</tr>")
    return "<table>" + "".join(html_rows) + "</table>"


def count_docx_embedded_objects(document: Any) -> int:
    count = 0
    try:
        for rel in document.part.rels.values():
            reltype = getattr(rel, "reltype", "")
            if "oleObject" in reltype or "package" in reltype:
                count += 1
    except Exception:
        return 0
    return count


def safe_int(value: Any) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def extract_slide_title(slide: Any) -> str | None:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        text = normalize_text(title_shape.text)
        return text or None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = normalize_text(shape.text)
            if text and len(text) <= 80:
                return text
    return None


def extract_slide_notes(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes = slide.notes_slide.notes_text_frame
        return normalize_text("\n".join(paragraph.text for paragraph in notes.paragraphs))
    except Exception:
        return ""
